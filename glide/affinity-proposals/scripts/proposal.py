from __future__ import annotations

import re
import urllib.parse
from dataclasses import dataclass
from datetime import date, timedelta
from pathlib import Path
from typing import Any

from scripts.extract import ProposalProfile
from scripts.seren_client import PublisherError


class SetupBlocked(RuntimeError):
    """Raised when an operator setup prerequisite blocks a dry-run."""


@dataclass
class ProposalTemplatePaths:
    offshore: Path
    onshore: Path


@dataclass
class ProposalArtifact:
    pptx_path: Path | None
    attachment_bytes: bytes
    file_name: str
    template_used: Path | None


def _presentation_class():
    from pptx import Presentation

    return Presentation


def _format_month_year(day: date) -> str:
    return day.strftime("%B %Y")


def _format_long_date(day: date) -> str:
    return day.strftime("%B %d, %Y").replace(" 0", " ")


def _safe_filename(name: str) -> str:
    value = re.sub(r"[^A-Za-z0-9._-]+", "_", name).strip("_")
    return value or "proposal"


def _replacement_map(profile: ProposalProfile, today: date) -> list[tuple[str, str]]:
    launch_date = today + timedelta(days=30)
    return [
        ("Secured Debt Investments", profile.client_name),
        ("Secured Debt", profile.client_name),
        ("CLIENT_NAME", profile.client_name),
        ("FUND_NAME", profile.fund_name),
        ("ADVISOR_NAME", profile.advisor_name),
        ("DESCRIPTION", profile.description),
        ("SEEKING", "\n".join(profile.seeking)),
        ("CURRENT_MONTH_YEAR", _format_month_year(today)),
        ("LAUNCH_DATE", _format_long_date(launch_date)),
    ]


def _replace_paragraph_text(paragraph: Any, replacements: list[tuple[str, str]]) -> None:
    original = "".join(run.text for run in paragraph.runs)
    if not original:
        return
    changed = original
    for old, new in replacements:
        changed = changed.replace(old, new)
    if changed == original:
        return
    if not paragraph.runs:
        paragraph.text = changed
        return
    paragraph.runs[0].text = changed
    for run in paragraph.runs[1:]:
        run.text = ""


def _iter_shapes(container: Any):
    """Yield every shape in `container`, descending into GROUP shapes.

    Slide-6's advisor/fund rectangles live inside nested groups (#980) — a
    flat `slide.shapes` walk would skip them.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in container.shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape)


def _iter_text_frames(slide: Any):
    for shape in _iter_shapes(slide):
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def write_proposal_deck(
    profile: ProposalProfile,
    *,
    templates: ProposalTemplatePaths,
    output_path: Path,
    today: date,
) -> ProposalArtifact:
    template = templates.onshore if profile.structure == "onshore" else templates.offshore
    prs = _presentation_class()(str(template))
    replacements = _replacement_map(profile, today)
    for index, slide in enumerate(prs.slides):
        if index == 9:
            continue
        for text_frame in _iter_text_frames(slide):
            for paragraph in text_frame.paragraphs:
                _replace_paragraph_text(paragraph, replacements)
    prs.core_properties.title = f"Glide - {profile.client_name} Proposal"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    file_name = f"{_safe_filename(profile.client_name)}_proposal_{today:%Y%m%d}.pptx"
    return ProposalArtifact(
        pptx_path=output_path,
        attachment_bytes=b"",
        file_name=file_name,
        template_used=template,
    )


def extract_text_by_slide(path: Path) -> list[str]:
    prs = _presentation_class()(str(path))
    slides: list[str] = []
    for slide in prs.slides:
        chunks: list[str] = []
        for text_frame in _iter_text_frames(slide):
            text = "\n".join(paragraph.text for paragraph in text_frame.paragraphs)
            if text:
                chunks.append(text)
        slides.append("\n".join(chunks))
    return slides


class SharePointRenderer:
    def __init__(self, gateway: Any, *, folder_name: str = "AI Proposals") -> None:
        self.gateway = gateway
        self.folder_name = folder_name

    def preflight(self) -> dict[str, Any]:
        try:
            return self.gateway.call_publisher(
                "microsoft-sharepoint", method="GET", path="/sites/root"
            )
        except PublisherError as exc:
            if self._is_oauth_error(exc):
                raise SetupBlocked(
                    "Microsoft OAuth connection required for the SharePoint render account. "
                    "Connect the render account to the microsoft-sharepoint publisher before dry-run."
                ) from exc
            raise

    @staticmethod
    def _is_oauth_error(exc: PublisherError) -> bool:
        return getattr(exc, "status", None) in (401, 403) or "oauth" in str(exc).lower()

    def _ensure_folder(self, drive_id: str) -> None:
        try:
            self.gateway.call_publisher(
                "microsoft-sharepoint",
                method="POST",
                path=f"/drives/{drive_id}/root/children",
                body={
                    "name": self.folder_name,
                    "folder": {},
                    "@microsoft.graph.conflictBehavior": "fail",
                },
            )
        except PublisherError as exc:
            # Folder already exists -> Graph returns 409 nameAlreadyExists.
            if getattr(exc, "status", None) != 409 and "alreadyexists" not in str(exc).lower():
                raise

    def render_pdf(self, pptx_path: Path) -> bytes:
        site = self.preflight()
        site_id = site.get("id") or site.get("site", {}).get("id")
        if not site_id:
            raise RuntimeError("SharePoint root site response missing id")
        drive = self.gateway.call_publisher(
            "microsoft-sharepoint", method="GET", path=f"/sites/{site_id}/drive"
        )
        drive_id = drive.get("id") or drive.get("drive", {}).get("id")
        if not drive_id:
            raise RuntimeError("SharePoint drive response missing id")

        self._ensure_folder(drive_id)
        folder = urllib.parse.quote(self.folder_name)
        name = urllib.parse.quote(pptx_path.name)
        upload = self.gateway.call_publisher(
            "microsoft-sharepoint",
            method="PUT",
            path=f"/drives/{drive_id}/root:/{folder}/{name}:/content",
            data=pptx_path.read_bytes(),
            content_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
        )
        item_id = upload.get("id") or upload.get("item", {}).get("id")
        if not item_id:
            raise RuntimeError("SharePoint upload response missing item id")
        pdf = self.gateway.call_publisher(
            "microsoft-sharepoint",
            method="GET",
            path=f"/drives/{drive_id}/items/{item_id}/content?format=pdf",
            response_format="binary",
        )
        return self._decode_pdf(pdf)

    @staticmethod
    def _decode_pdf(raw: Any) -> bytes:
        # The gateway returns binary downloads base64-encoded in its JSON
        # envelope; GatewayClient.call_publisher(response_format="binary")
        # decodes that to the raw bytes (seren-core #182). Render therefore
        # receives real PDF bytes here.
        pdf_bytes = raw.encode("latin1", "replace") if isinstance(raw, str) else bytes(raw)
        if pdf_bytes.startswith(b"%PDF"):
            return pdf_bytes
        raise RuntimeError("SharePoint render did not return PDF bytes")


class ProposalService:
    def __init__(
        self,
        *,
        templates: ProposalTemplatePaths,
        output_dir: Path,
    ) -> None:
        self.templates = templates
        self.output_dir = output_dir

    def build(self, profile: ProposalProfile, today: date) -> ProposalArtifact:
        # The skill ships the editable .pptx as the email attachment (#980);
        # there is no PDF render step anymore.
        self.output_dir.mkdir(parents=True, exist_ok=True)
        file_name = f"{_safe_filename(profile.client_name)}_proposal_{today:%Y%m%d}.pptx"
        out_path = self.output_dir / file_name
        artifact = write_proposal_deck(
            profile,
            templates=self.templates,
            output_path=out_path,
            today=today,
        )
        pptx_bytes = out_path.read_bytes()
        return ProposalArtifact(
            pptx_path=out_path,
            attachment_bytes=pptx_bytes,
            file_name=artifact.file_name,
            template_used=artifact.template_used,
        )
